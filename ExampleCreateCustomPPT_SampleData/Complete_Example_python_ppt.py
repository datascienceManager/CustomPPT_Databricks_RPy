"""
============================================================================
COMPLETE EXAMPLE: Corporate Quarterly Report - PowerPoint Generator
============================================================================
This creates a professional PowerPoint presentation using charts from R
Theme: Q4 2024 Business Performance Report
Works with or without a template!
============================================================================
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

print("=" * 80)
print("CORPORATE QUARTERLY REPORT - POWERPOINT GENERATOR")
print("=" * 80)
print()

# ============================================================================
# CONFIGURATION
# ============================================================================

# Template path (optional - works without template too!)
TEMPLATE_PATH = "/dbfs/FileStore/templates/company_template.pptx"
USE_TEMPLATE = os.path.exists(TEMPLATE_PATH)

# Chart paths
CHART_DIR = "/dbfs/FileStore/example_charts"
DATA_DIR = "/dbfs/FileStore/example_data"

# Output path
OUTPUT_DIR = "/dbfs/FileStore/presentations"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ============================================================================
# CORPORATE COLOR SCHEME
# ============================================================================

class CorporateColors:
    """Corporate color palette"""
    PRIMARY = RGBColor(31, 71, 136)      # Deep Blue #1f4788
    SECONDARY = RGBColor(113, 165, 222)  # Light Blue #71a5de
    ACCENT1 = RGBColor(243, 156, 18)     # Orange #f39c12
    ACCENT2 = RGBColor(39, 174, 96)      # Green #27ae60
    ACCENT3 = RGBColor(231, 76, 60)      # Red #e74c3c
    NEUTRAL = RGBColor(149, 165, 166)    # Gray #95a5a6
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(44, 62, 80)
    LIGHT_TEXT = RGBColor(127, 140, 141)
    BACKGROUND = RGBColor(236, 240, 241)

colors = CorporateColors()

# ============================================================================
# LOAD DATA
# ============================================================================

print("1. Loading data from CSV files...")

try:
    summary_metrics = pd.read_csv(f"{DATA_DIR}/summary_metrics.csv")
    monthly_sales = pd.read_csv(f"{DATA_DIR}/monthly_sales.csv")
    product_performance = pd.read_csv(f"{DATA_DIR}/product_performance.csv")
    regional_sales = pd.read_csv(f"{DATA_DIR}/regional_sales.csv")
    customer_segments = pd.read_csv(f"{DATA_DIR}/customer_segments.csv")
    department_budget = pd.read_csv(f"{DATA_DIR}/department_budget.csv")
    
    print("   ✓ All data files loaded successfully")
    print()
except Exception as e:
    print(f"   ✗ Error loading data: {e}")
    print("   Make sure to run the R script first!")
    exit(1)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0] if USE_TEMPLATE else prs.slide_layouts[6])
    
    if not USE_TEMPLATE:
        # Create title slide from scratch
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors.PRIMARY
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = colors.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = colors.SECONDARY
        p.alignment = PP_ALIGN.CENTER
    else:
        # Use template placeholders
        try:
            slide.placeholders[0].text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
        except:
            pass
    
    return slide

def add_content_slide_with_chart(prs, title, chart_path, layout_idx=1):
    """Add a slide with title and chart"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx] if USE_TEMPLATE else prs.slide_layouts[6])
    
    if not USE_TEMPLATE:
        # Create from scratch
        # Title bar
        title_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(0.8)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = colors.PRIMARY
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = colors.WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Chart
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5),   # left
                Inches(1.3),   # top (below title)
                width=Inches(9)  # width
            )
            print(f"   ✓ Chart added: {title}")
        else:
            print(f"   ✗ Chart not found: {chart_path}")
    else:
        # Use template
        try:
            slide.placeholders[0].text = title
            
            # Find content placeholder and replace with chart
            if len(slide.placeholders) > 1:
                content_ph = slide.placeholders[1]
                left = content_ph.left.inches if content_ph.left > 0 else 0.5
                top = content_ph.top.inches if content_ph.top > 0 else 1.5
                width = content_ph.width.inches
                
                # Remove placeholder
                sp = content_ph.element
                sp.getparent().remove(sp)
                
                # Add chart
                if os.path.exists(chart_path):
                    slide.shapes.add_picture(chart_path, Inches(left), Inches(top), width=Inches(width))
                    print(f"   ✓ Chart added: {title}")
            else:
                # Fallback positioning
                if os.path.exists(chart_path):
                    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(9))
        except Exception as e:
            print(f"   ⚠ Using fallback positioning: {e}")
            if os.path.exists(chart_path):
                slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(9))
    
    return slide

def add_metrics_slide(prs, title, metrics_df):
    """Create executive summary slide with metric cards"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors.PRIMARY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    
    # Metric cards (2x2 grid)
    for idx, (_, metric) in enumerate(metrics_df.iterrows()):
        col = idx % 2
        row = idx // 2
        
        left = 0.5 + (col * 4.75)
        top = 1.3 + (row * 1.6)
        
        # Card background
        card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(left), Inches(top), Inches(4.25), Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors.BACKGROUND
        card.line.color.rgb = colors.PRIMARY
        card.line.width = Pt(2)
        
        # Metric name
        name_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.15),
            Inches(4), Inches(0.3)
        )
        name_box.text = metric['Metric']
        p = name_box.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors.DARK_TEXT
        
        # Value
        value_text = f"${metric['Value']}{metric['Unit']}" if metric['Unit'] in ['M', 'K'] else f"{metric['Value']}{metric['Unit']}"
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.5),
            Inches(2.5), Inches(0.5)
        )
        value_box.text = value_text
        p = value_box.text_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors.PRIMARY
        
        # Change indicator
        change_text = f"▲ {metric['Change']}%" if metric['Change'] > 0 else f"▼ {abs(metric['Change'])}%"
        change_box = slide.shapes.add_textbox(
            Inches(left + 3), Inches(top + 0.6),
            Inches(1), Inches(0.4)
        )
        change_box.text = change_text
        p = change_box.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors.ACCENT2 if metric['Change'] > 0 else colors.ACCENT3
        p.alignment = PP_ALIGN.RIGHT
    
    print(f"   ✓ Metrics slide created")
    return slide

def add_table_slide(prs, title, df, columns_to_show):
    """Create a slide with a data table"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors.PRIMARY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    
    # Table
    rows = len(df) + 1  # +1 for header
    cols = len(columns_to_show)
    
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(1), Inches(1.3),
        Inches(8), Inches(3.5)
    )
    table = table_shape.table
    
    # Header row
    for col_idx, col_name in enumerate(columns_to_show):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors.PRIMARY
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = colors.WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, (_, row) in enumerate(df.iterrows(), start=1):
        for col_idx, col_name in enumerate(columns_to_show):
            cell = table.cell(row_idx, col_idx)
            
            # Format value
            value = row[col_name]
            if isinstance(value, float):
                if value > 100:
                    cell.text = f"${value:.1f}M"
                else:
                    cell.text = f"{value:.1f}"
            else:
                cell.text = str(value)
            
            # Styling
            bg_color = colors.WHITE if row_idx % 2 == 1 else colors.BACKGROUND
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
    
    print(f"   ✓ Table slide created: {title}")
    return slide

def add_closing_slide(prs):
    """Create closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors.PRIMARY
    
    # Thank you text
    thank_you = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    thank_you.text = "THANK YOU"
    p = thank_you.text_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Questions text
    questions = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.5))
    questions.text = "Questions & Discussion"
    p = questions.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = colors.SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================================
# CREATE PRESENTATION
# ============================================================================

print("2. Creating PowerPoint presentation...")
print()

if USE_TEMPLATE:
    prs = Presentation(TEMPLATE_PATH)
    print(f"   ✓ Using template: {TEMPLATE_PATH}")
else:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    print("   ✓ Creating presentation without template (blank slides)")

print()
print("3. Adding slides...")
print()

# Slide 1: Title Slide
add_title_slide(
    prs,
    "Q4 2024 BUSINESS PERFORMANCE",
    "Quarterly Executive Report | December 2024"
)
print("   ✓ Slide 1: Title slide")

# Slide 2: Executive Summary
add_metrics_slide(
    prs,
    "EXECUTIVE SUMMARY",
    summary_metrics
)
print("   ✓ Slide 2: Executive summary with metrics")

# Slide 3: Revenue Trend
add_content_slide_with_chart(
    prs,
    "REVENUE PERFORMANCE",
    f"{CHART_DIR}/01_revenue_trend.png"
)

# Slide 4: Product Performance
add_content_slide_with_chart(
    prs,
    "PRODUCT PERFORMANCE",
    f"{CHART_DIR}/02_product_performance.png"
)

# Slide 5: Regional Distribution
add_content_slide_with_chart(
    prs,
    "REGIONAL DISTRIBUTION",
    f"{CHART_DIR}/03_regional_distribution.png"
)

# Slide 6: Customer Segments
add_content_slide_with_chart(
    prs,
    "CUSTOMER SEGMENTS",
    f"{CHART_DIR}/04_customer_segments.png"
)

# Slide 7: Department Budget
add_content_slide_with_chart(
    prs,
    "DEPARTMENT BUDGET",
    f"{CHART_DIR}/05_department_budget.png"
)

# Slide 8: Growth Rates
add_content_slide_with_chart(
    prs,
    "PRODUCT GROWTH RATES",
    f"{CHART_DIR}/06_growth_rates.png"
)

# Slide 9: Product Performance Table
add_table_slide(
    prs,
    "DETAILED PRODUCT METRICS",
    product_performance,
    ['Product', 'Q3_Sales', 'Q4_Sales', 'Growth']
)

# Slide 10: Closing
add_closing_slide(prs)
print("   ✓ Slide 10: Closing slide")

# ============================================================================
# SAVE PRESENTATION
# ============================================================================

print()
print("4. Saving presentation...")

output_path = f"{OUTPUT_DIR}/Q4_2024_Business_Report.pptx"
prs.save(output_path)

print(f"   ✓ Presentation saved: {output_path}")
print()

# ============================================================================
# SUMMARY
# ============================================================================

print("=" * 80)
print("POWERPOINT CREATION COMPLETE!")
print("=" * 80)
print()
print("✅ Presentation Details:")
print(f"   - Total Slides: 10")
print(f"   - Template Used: {'Yes' if USE_TEMPLATE else 'No (created from scratch)'}")
print(f"   - File Size: {os.path.getsize(output_path) / 1024:.1f} KB")
print()
print("📊 Content Summary:")
print("   - 1 Title slide")
print("   - 1 Executive summary (4 metric cards)")
print("   - 6 Chart slides (revenue, products, regions, segments, budget, growth)")
print("   - 1 Data table slide")
print("   - 1 Closing slide")
print()
print("📥 Download Location:")
print(f"   {output_path}")
print()
print("=" * 80)
print()
print("💡 POSITIONING USED:")
print("   - Charts: 0.5\" from left, 1.3\" from top, 9\" wide")
print("   - Title bar: Full width, 0.8\" tall")
print("   - Metric cards: 2x2 grid with 0.5\" margins")
print("=" * 80)
