"""
============================================================================
SAMPLE TEMPLATE GENERATOR
============================================================================
This creates a simple PowerPoint template that you can customize
Use this as a starting point for your own corporate template
============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

print("Creating sample PowerPoint template...")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Define colors (customize these!)
PRIMARY_COLOR = RGBColor(31, 71, 136)      # Deep Blue
SECONDARY_COLOR = RGBColor(113, 165, 222)   # Light Blue
ACCENT_COLOR = RGBColor(243, 156, 18)       # Orange

# ============================================================================
# LAYOUT 1: TITLE SLIDE
# ============================================================================

title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)

# Customize title slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PRIMARY_COLOR

# Add company logo placeholder (you can replace this)
logo_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.2), Inches(1.3), Inches(0.4))
logo_box.text = "YOUR LOGO"
p = logo_box.text_frame.paragraphs[0]
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = SECONDARY_COLOR
p.alignment = PP_ALIGN.CENTER

# Title
try:
    slide.placeholders[0].text = "YOUR COMPANY NAME"
    title_frame = slide.placeholders[0].text_frame
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
except:
    pass

# Subtitle
try:
    slide.placeholders[1].text = "Presentation Title | Date"
    subtitle_frame = slide.placeholders[1].text_frame
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = SECONDARY_COLOR
except:
    pass

# ============================================================================
# LAYOUT 2: CONTENT SLIDE EXAMPLE
# ============================================================================

content_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(content_layout)

# Add branded header
header = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(0), Inches(10), Inches(0.7)
)
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY_COLOR
header.line.fill.background()

# Add small logo to header
logo_box = slide.shapes.add_textbox(Inches(9), Inches(0.15), Inches(0.8), Inches(0.4))
logo_box.text = "LOGO"
p = logo_box.text_frame.paragraphs[0]
p.font.size = Pt(10)
p.font.color.rgb = SECONDARY_COLOR
p.alignment = PP_ALIGN.CENTER

# Customize placeholders
try:
    slide.placeholders[0].text = "Content Slide Example"
    title_frame = slide.placeholders[0].text_frame
    title_frame.clear()
    
    # Create title in custom position
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(8), Inches(0.5))
    title_box.text = "Slide Title Goes Here"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
except:
    pass

# ============================================================================
# SAVE TEMPLATE
# ============================================================================

output_path = "/dbfs/FileStore/templates/sample_template.pptx"
import os
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)

print(f"✅ Sample template created: {output_path}")
print()
print("📝 To customize:")
print("   1. Download the template")
print("   2. Open in PowerPoint")
print("   3. Add your logo, colors, fonts")
print("   4. Save and re-upload")
print("   5. Use in your presentations!")
