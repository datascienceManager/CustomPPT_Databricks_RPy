"""
============================================================================
DATABRICKS PYTHON SCRIPT: POWERPOINT PRESENTATION BUILDER WITH TEMPLATE
============================================================================
Purpose: Read data from R Spark temp views and create PowerPoint using YOUR template
Input: 
  - Spark temp views + PNG charts from /dbfs/FileStore/charts/
  - YOUR TEMPLATE: /dbfs/FileStore/templates/your_template.pptx
Output: PowerPoint presentation in /dbfs/FileStore/presentations/
============================================================================
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime
import os

print("=" * 80)
print("SPORTS VIEWING ANALYTICS - POWERPOINT BUILDER (TEMPLATE MODE)")
print("=" * 80)
print()

# ============================================================================
# CONFIGURATION: SET YOUR TEMPLATE PATH
# ============================================================================

# IMPORTANT: Upload your template to Databricks first, then set the path here
TEMPLATE_PATH = "/dbfs/FileStore/templates/your_template.pptx"

# If template doesn't exist, use blank presentation
USE_TEMPLATE = os.path.exists(TEMPLATE_PATH)

if USE_TEMPLATE:
    print(f"✓ Using template: {TEMPLATE_PATH}")
else:
    print(f"⚠ Template not found at: {TEMPLATE_PATH}")
    print("  Creating presentation without template (blank slides)")
    print("  To use a template:")
    print("    1. Upload your .pptx template to Databricks")
    print("    2. Update TEMPLATE_PATH variable above")
print()

# ============================================================================
# SECTION 1: LOAD DATA FROM R SPARK TEMP VIEWS
# ============================================================================
print("1. Loading data from R Spark temp views...")

try:
    summary_data = spark.sql("SELECT * FROM summary_data").toPandas()
    monthly_data = spark.sql("SELECT * FROM monthly_data").toPandas()
    sports_data = spark.sql("SELECT * FROM sports_data").toPandas()
    device_data = spark.sql("SELECT * FROM device_data").toPandas()
    
    # Convert summary to dictionary
    summary = summary_data.iloc[0].to_dict()
    
    print("   ✓ Summary data loaded")
    print("   ✓ Monthly data loaded:", len(monthly_data), "rows")
    print("   ✓ Sports data loaded:", len(sports_data), "rows")
    print("   ✓ Device data loaded:", len(device_data), "rows")
    print()
    
except Exception as e:
    print(f"   ✗ Error loading data: {e}")
    print("   Make sure to run the R script first!")
    raise

# ============================================================================
# SECTION 2: GENERATE INSIGHTS & RECOMMENDATIONS
# ============================================================================
print("2. Generating insights and recommendations...")

def generate_key_findings(summary, monthly_data, sports_data, device_data):
    """Generate key findings from the data"""
    findings = []
    
    findings.append(
        f"Total viewing time reached {int(summary['TotalViewingMinutes']):,} minutes "
        f"across {int(summary['UniqueViewers']):,} unique viewers"
    )
    
    top_comp = monthly_data.groupby('Competition')['ViewingMinutes'].sum().idxmax()
    top_comp_views = int(monthly_data.groupby('Competition')['ViewingMinutes'].sum().max())
    findings.append(
        f"{top_comp} led all competitions with {top_comp_views:,} total viewing minutes"
    )
    
    avg_engagement = summary['AvgMinutesPerViewer']
    findings.append(
        f"Average viewer engagement of {avg_engagement:.1f} minutes per viewer "
        f"indicates strong content retention"
    )
    
    top_sport = sports_data.nlargest(1, 'ViewingMinutes')['Sport'].values[0]
    findings.append(f"{top_sport} dominated sports viewership across all categories")
    
    top_device = device_data.nlargest(1, 'UniqueViewers')['Device'].values[0]
    findings.append(
        f"{top_device} was the preferred viewing device, capturing the largest audience share"
    )
    
    monthly_totals = monthly_data.groupby('YearMonth')['ViewingMinutes'].sum()
    first_half = monthly_totals.iloc[:6].mean()
    second_half = monthly_totals.iloc[6:].mean()
    growth = ((second_half - first_half) / first_half) * 100
    findings.append(
        f"Viewing minutes {'increased' if growth > 0 else 'decreased'} by "
        f"{abs(growth):.1f}% in the second half of the year"
    )
    
    return findings

def generate_recommendations(monthly_data, sports_data, device_data):
    """Generate strategic recommendations"""
    recommendations = []
    
    top_3_comps = monthly_data.groupby('Competition')['UniqueViewers'].sum().nlargest(3).index.tolist()
    recommendations.append(
        f"Focus marketing efforts on top-performing competitions: {', '.join(top_3_comps)}"
    )
    
    mobile_users = device_data[device_data['Device'] == 'Mobile']['UniqueViewers'].values
    tv_users = device_data[device_data['Device'] == 'TV']['UniqueViewers'].values
    
    if len(mobile_users) > 0 and len(tv_users) > 0:
        if mobile_users[0] > tv_users[0]:
            recommendations.append(
                "Prioritize mobile app enhancements and responsive design given strong mobile adoption"
            )
        else:
            recommendations.append(
                "Optimize TV viewing experience with enhanced picture quality and interactive features"
            )
    
    low_engagement_sports = sports_data.nsmallest(2, 'MinutesPerViewer')['Sport'].tolist()
    if len(low_engagement_sports) >= 2:
        recommendations.append(
            f"Improve content quality and promotion for {' and '.join(low_engagement_sports)} "
            f"to boost engagement"
        )
    
    recommendations.append(
        "Develop targeted campaigns for Q4 when football competitions peak in viewership"
    )
    
    recommendations.append(
        "Implement seamless cross-device experiences to support viewers who switch between TV and mobile"
    )
    
    recommendations.append(
        "Launch loyalty programs to convert casual viewers into regular subscribers"
    )
    
    return recommendations

findings = generate_key_findings(summary, monthly_data, sports_data, device_data)
recommendations = generate_recommendations(monthly_data, sports_data, device_data)

print("   ✓ Generated", len(findings), "key findings")
print("   ✓ Generated", len(recommendations), "strategic recommendations")
print()

# ============================================================================
# SECTION 3: HELPER FUNCTIONS FOR TEMPLATE HANDLING
# ============================================================================

def get_slide_layout_info(prs):
    """Display information about available slide layouts in the template"""
    print("\n   📋 Available Slide Layouts in Template:")
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"      Layout {idx}: {layout.name}")
        print(f"         Placeholders: {len(layout.placeholders)}")
        for pidx, placeholder in enumerate(layout.placeholders):
            print(f"           [{pidx}] {placeholder.name} ({placeholder.placeholder_format.type})")
    print()

def add_content_to_placeholder(slide, placeholder_idx, content, font_size=16, bold=False, color=None):
    """Add text content to a specific placeholder"""
    try:
        placeholder = slide.placeholders[placeholder_idx]
        text_frame = placeholder.text_frame
        text_frame.clear()
        text_frame.text = content
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            if bold:
                paragraph.font.bold = True
            if color:
                paragraph.font.color.rgb = color
        
        return True
    except (KeyError, IndexError):
        print(f"      Warning: Placeholder {placeholder_idx} not found")
        return False

def add_bullet_points_to_placeholder(slide, placeholder_idx, items, font_size=16):
    """Add bullet points to a placeholder"""
    try:
        placeholder = slide.placeholders[placeholder_idx]
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        for idx, item in enumerate(items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(font_size)
            p.space_after = Pt(12)
        
        return True
    except (KeyError, IndexError):
        print(f"      Warning: Placeholder {placeholder_idx} not found")
        return False

def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add image to slide at specified position"""
    if os.path.exists(image_path):
        if width and height:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                    width=Inches(width), height=Inches(height))
        elif width:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
        else:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top))
        return True
    else:
        print(f"      Warning: Image not found: {image_path}")
        return False

# ============================================================================
# SECTION 4: CREATE POWERPOINT PRESENTATION
# ============================================================================
print("3. Creating PowerPoint presentation...")

def create_presentation_with_template(summary, sports_data, device_data, findings, recommendations):
    """Create the PowerPoint presentation using template"""
    
    # Load template or create blank presentation
    if USE_TEMPLATE:
        prs = Presentation(TEMPLATE_PATH)
        print("   ✓ Template loaded successfully")
        
        # Show available layouts (helpful for first-time setup)
        # Uncomment the line below to see what layouts are available
        # get_slide_layout_info(prs)
    else:
        prs = Presentation()
        print("   ✓ Created blank presentation")
    
    print("   Creating slides...")
    
    # ========================================================================
    # CONFIGURATION: CUSTOMIZE THESE LAYOUT INDICES FOR YOUR TEMPLATE
    # ========================================================================
    # These indices depend on your template structure
    # To find the right indices, uncomment get_slide_layout_info(prs) above
    
    LAYOUT_TITLE = 0           # Title slide layout index
    LAYOUT_SECTION = 1         # Section header layout (if available)
    LAYOUT_TITLE_CONTENT = 1   # Title + content layout
    LAYOUT_TITLE_ONLY = 5      # Title only layout (for charts)
    LAYOUT_BLANK = 6           # Blank layout
    
    # Try to use template layouts, fall back to indices 0-6 if not available
    def get_layout(index, fallback=0):
        try:
            return prs.slide_layouts[index]
        except IndexError:
            print(f"      Warning: Layout {index} not found, using layout {fallback}")
            return prs.slide_layouts[fallback]
    
    # ========================================================================
    # SLIDE 1: TITLE SLIDE (Using Template Layout)
    # ========================================================================
    print("      - Slide 1: Title")
    slide1 = prs.slides.add_slide(get_layout(LAYOUT_TITLE))
    
    # Method 1: Use template placeholders (PREFERRED)
    # Find title and subtitle placeholders
    try:
        # Common placeholder indices for title slides
        # Title is usually placeholder 0, subtitle is usually placeholder 1
        add_content_to_placeholder(slide1, 0, "SPORTS VIEWING ANALYTICS REPORT 2025", 
                                  font_size=44, bold=True)
        add_content_to_placeholder(slide1, 1, 
                                  "Comprehensive Analysis of Viewing Trends, Engagement & Recommendations",
                                  font_size=18)
    except:
        # Method 2: Add text manually if placeholders don't work
        title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "SPORTS VIEWING ANALYTICS REPORT 2025"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # ========================================================================
    print("      - Slide 2: Executive Summary")
    slide2 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_CONTENT))
    
    add_content_to_placeholder(slide2, 0, "EXECUTIVE SUMMARY", font_size=36, bold=True)
    
    # Create metrics text
    metrics_text = f"""
📊 Total Viewing Minutes: {int(summary['TotalViewingMinutes']):,}

👥 Unique Viewers: {int(summary['UniqueViewers']):,}

⏱️ Average Minutes per Viewer: {summary['AvgMinutesPerViewer']:.1f}

🎬 Total Assets: {int(summary['TotalAssets']):,}
    """.strip()
    
    add_content_to_placeholder(slide2, 1, metrics_text, font_size=20)
    
    # ========================================================================
    # SLIDE 3: KEY FINDINGS
    # ========================================================================
    print("      - Slide 3: Key Findings")
    slide3 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_CONTENT))
    
    add_content_to_placeholder(slide3, 0, "KEY FINDINGS", font_size=36, bold=True)
    add_bullet_points_to_placeholder(slide3, 1, findings, font_size=16)
    
    # ========================================================================
    # SLIDE 4: SPORTS DISTRIBUTION CHART
    # ========================================================================
    print("      - Slide 4: Sports Distribution Chart")
    slide4 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_ONLY))
    
    add_content_to_placeholder(slide4, 0, "VIEWERSHIP BY SPORT", font_size=36, bold=True)
    add_image_to_slide(slide4, "/dbfs/FileStore/charts/sports_pie.png", 
                      left=1, top=1.5, width=8)
    
    # ========================================================================
    # SLIDE 5: DEVICE VIEWERSHIP CHART
    # ========================================================================
    print("      - Slide 5: Device Viewership Chart")
    slide5 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_ONLY))
    
    add_content_to_placeholder(slide5, 0, "VIEWERSHIP BY DEVICE", font_size=36, bold=True)
    add_image_to_slide(slide5, "/dbfs/FileStore/charts/device_lollipop.png", 
                      left=1, top=1.5, width=8)
    
    # ========================================================================
    # SLIDE 6: COMPETITION TRENDS CHART
    # ========================================================================
    print("      - Slide 6: Competition Trends Chart")
    slide6 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_ONLY))
    
    add_content_to_placeholder(slide6, 0, "COMPETITION TRENDS OVER TIME", 
                              font_size=32, bold=True)
    add_image_to_slide(slide6, "/dbfs/FileStore/charts/competition_line.png", 
                      left=0.5, top=1.5, width=9)
    
    # ========================================================================
    # SLIDE 7: TOP COMPETITIONS CHART
    # ========================================================================
    print("      - Slide 7: Top Competitions Chart")
    slide7 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_ONLY))
    
    add_content_to_placeholder(slide7, 0, "TOP PERFORMING COMPETITIONS", 
                              font_size=36, bold=True)
    add_image_to_slide(slide7, "/dbfs/FileStore/charts/competition_bar.png", 
                      left=1, top=1.5, width=8)
    
    # ========================================================================
    # SLIDE 8: GT/FLEXTABLE (if available)
    # ========================================================================
    if os.path.exists("/dbfs/FileStore/charts/sports_gt_table.png"):
        print("      - Slide 8: Sports Performance Table")
        slide8 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_ONLY))
        
        add_content_to_placeholder(slide8, 0, "DETAILED SPORTS METRICS", 
                                  font_size=36, bold=True)
        add_image_to_slide(slide8, "/dbfs/FileStore/charts/sports_gt_table.png", 
                          left=0.5, top=1.5, width=9)
    else:
        # Fallback: Create table using python-pptx if GT table image doesn't exist
        print("      - Slide 8: Sports Metrics Table (python-pptx)")
        slide8 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_CONTENT))
        
        add_content_to_placeholder(slide8, 0, "DETAILED SPORTS METRICS", 
                                  font_size=36, bold=True)
        
        # Create table in placeholder or manually
        try:
            # Remove content placeholder and add table
            content_placeholder = slide8.placeholders[1]
            sp = content_placeholder.element
            sp.getparent().remove(sp)
        except:
            pass
        
        # Add table manually
        rows = len(sports_data) + 1
        cols = 4
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(3)
        
        table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Headers
        headers = ["Sport", "Viewing Minutes", "Unique Viewers", "Minutes/Viewer"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Data
        for row_idx, (_, row) in enumerate(sports_data.iterrows(), start=1):
            table.cell(row_idx, 0).text = str(row['Sport'])
            table.cell(row_idx, 1).text = f"{int(row['ViewingMinutes']):,}"
            table.cell(row_idx, 2).text = f"{int(row['UniqueViewers']):,}"
            table.cell(row_idx, 3).text = f"{row['MinutesPerViewer']:.2f}"
            
            for col_idx in range(4):
                cell = table.cell(row_idx, col_idx)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SLIDE 9: RECOMMENDATIONS
    # ========================================================================
    print("      - Slide 9: Strategic Recommendations")
    slide9 = prs.slides.add_slide(get_layout(LAYOUT_TITLE_CONTENT))
    
    add_content_to_placeholder(slide9, 0, "STRATEGIC RECOMMENDATIONS", 
                              font_size=36, bold=True)
    
    # Format recommendations with numbers
    numbered_recs = [f"{i+1}. {rec}" for i, rec in enumerate(recommendations)]
    add_bullet_points_to_placeholder(slide9, 1, numbered_recs, font_size=16)
    
    # ========================================================================
    # SLIDE 10: CLOSING SLIDE (Using Template)
    # ========================================================================
    print("      - Slide 10: Closing")
    slide10 = prs.slides.add_slide(get_layout(LAYOUT_TITLE))
    
    try:
        add_content_to_placeholder(slide10, 0, "THANK YOU", font_size=56, bold=True)
        add_content_to_placeholder(slide10, 1, "Questions & Discussion", font_size=24)
    except:
        # Manual fallback
        title_box = slide10.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "THANK YOU"
        title_frame.paragraphs[0].font.size = Pt(56)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========================================================================
    # SAVE PRESENTATION
    # ========================================================================
    print("   Saving presentation...")
    
    os.makedirs('/dbfs/FileStore/presentations', exist_ok=True)
    output_path = '/dbfs/FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx'
    prs.save(output_path)
    
    return output_path

# Create the presentation
try:
    output_file = create_presentation_with_template(summary, sports_data, device_data, 
                                                    findings, recommendations)
    print("   ✓ Presentation saved successfully")
    print()
except Exception as e:
    print(f"   ✗ Error creating presentation: {e}")
    import traceback
    traceback.print_exc()
    raise

# ============================================================================
# SECTION 5: COMPLETION SUMMARY
# ============================================================================
print("=" * 80)
print("POWERPOINT CREATION COMPLETE!")
print("=" * 80)
print()
print("✅ Presentation Details:")
print(f"   - Template Used: {'Yes' if USE_TEMPLATE else 'No (blank slides)'}")
print(f"   - Total Slides: 10")
print(f"   - File Location: {output_file}")
print(f"   - File Size: {os.path.getsize(output_file) / 1024:.1f} KB")
print()
print("📊 Content Summary:")
print(f"   - 1 Title slide")
print(f"   - 1 Executive summary")
print(f"   - 1 Key findings slide with {len(findings)} insights")
print(f"   - 4 Chart slides")
print(f"   - 1 Data table slide")
print(f"   - 1 Recommendations slide with {len(recommendations)} actions")
print(f"   - 1 Closing slide")
print()
print("📥 Download Instructions:")
print("   In Databricks, download from:")
print("   /FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx")
print()
print("=" * 80)

# Display download link
displayHTML(f"""
<div style="background-color: #E8F4F5; padding: 25px; border-radius: 10px; border: 2px solid #028090;">
    <h2 style="color: #028090; margin-top: 0;">✅ Presentation Ready!</h2>
    <p style="color: #1F2937; font-size: 16px;">
        Your sports viewing analytics presentation has been created 
        {'<strong>using your custom template</strong>' if USE_TEMPLATE else '<strong>with blank slides</strong>'}.
    </p>
    <p style="color: #6B7280; font-size: 14px; margin-bottom: 20px;">
        Location: <code>/FileStore/presentations/Sports_Viewing_Analytics_Report_2025.pptx</code>
    </p>
    <a href="/files/presentations/Sports_Viewing_Analytics_Report_2025.pptx" download>
        <button style="background-color: #028090; color: white; padding: 12px 24px; 
                       border: none; border-radius: 5px; cursor: pointer; font-size: 16px;
                       font-weight: bold; box-shadow: 0 2px 4px rgba(0,0,0,0.2);">
            📥 Download PowerPoint Presentation
        </button>
    </a>
    
    {'<div style="margin-top: 20px; padding: 15px; background-color: #FFF3CD; border-radius: 5px; border-left: 4px solid #FFC107;"><strong>⚠️ Note:</strong> Template was not found. To use your custom template:<ol style="margin: 10px 0 0 0; padding-left: 20px;"><li>Upload your .pptx template to <code>/FileStore/templates/your_template.pptx</code></li><li>Update the <code>TEMPLATE_PATH</code> variable in the script</li><li>Re-run this script</li></ol></div>' if not USE_TEMPLATE else '<div style="margin-top: 20px; padding: 15px; background-color: #D1F2EB; border-radius: 5px; border-left: 4px solid #02C39A;"><strong>✓ Success:</strong> Your custom template was used successfully!</div>'}
    
    <div style="margin-top: 20px; padding: 15px; background-color: white; border-radius: 5px;">
        <h4 style="color: #028090; margin-top: 0;">Presentation Contains:</h4>
        <ul style="color: #1F2937; line-height: 1.8;">
            <li><strong>{len(findings)}</strong> automated key findings</li>
            <li><strong>{len(recommendations)}</strong> strategic recommendations</li>
            <li><strong>4</strong> professional charts from R</li>
            <li><strong>1</strong> detailed data table</li>
            <li><strong>10</strong> total slides</li>
        </ul>
    </div>
</div>
""")
