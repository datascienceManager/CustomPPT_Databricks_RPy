"""
============================================================================
PRACTICAL POSITIONING EXAMPLES - Copy & Use These Patterns
============================================================================
This script shows you EXACTLY how to position charts in common scenarios.
Just copy the pattern that matches your template!
============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# ============================================================================
# PATTERN 1: CHART FILLING A CONTENT PLACEHOLDER
# ============================================================================
# Use when: Your template has a "Title and Content" layout
# The chart should fill the content area

def pattern_1_fill_placeholder(prs, layout_idx, title_text, chart_path):
    """Chart fills the content placeholder exactly"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add title (usually placeholder 0)
    try:
        title = slide.placeholders[0]
        title.text = title_text
    except:
        pass
    
    # Get content placeholder dimensions (usually placeholder 1)
    try:
        content_ph = slide.placeholders[1]
        
        # Get exact position and size
        left = content_ph.left.inches if content_ph.left > 0 else 0
        top = content_ph.top.inches if content_ph.top > 0 else 0
        width = content_ph.width.inches
        height = content_ph.height.inches
        
        # Remove placeholder so it doesn't show
        sp = content_ph.element
        sp.getparent().remove(sp)
        
        # Add chart in exact same space
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
            print(f"   ✓ Pattern 1: Chart added to placeholder area")
        
    except Exception as e:
        print(f"   ✗ Pattern 1 failed: {e}")
        # Fallback: use manual positioning
        pattern_2_below_title(prs, layout_idx, title_text, chart_path)

# ============================================================================
# PATTERN 2: CHART BELOW TITLE (MOST COMMON)
# ============================================================================
# Use when: Simple layout with title at top, chart below
# Works with most templates

def pattern_2_below_title(prs, layout_idx, title_text, chart_path):
    """Chart positioned below title with standard margins"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add title
    try:
        title = slide.placeholders[0]
        title.text = title_text
        
        # Find where title ends
        title_bottom = title.top.inches + title.height.inches
        
        # Position chart below title
        chart_left = 1.0      # 1 inch from left edge
        chart_top = title_bottom + 0.3   # 0.3 inch below title
        chart_width = 8.0     # 8 inches wide (fits in 10" slide with margins)
        
    except:
        # If can't find title, use absolute position
        chart_left = 1.0
        chart_top = 1.5
        chart_width = 8.0
    
    # Add chart
    if os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            Inches(chart_left),
            Inches(chart_top),
            width=Inches(chart_width)
        )
        print(f"   ✓ Pattern 2: Chart positioned below title")
    else:
        print(f"   ✗ Chart not found: {chart_path}")

# ============================================================================
# PATTERN 3: CENTERED CHART (BLANK LAYOUT)
# ============================================================================
# Use when: Using blank layout, want chart centered

def pattern_3_centered(prs, layout_idx, title_text, chart_path):
    """Chart centered on slide with title at top"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add title manually (since blank layout has no placeholders)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6)
    )
    title_box.text_frame.text = title_text
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Center the chart
    chart_width = 8.0
    chart_height = 4.0  # Approximate
    
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    chart_left = (slide_width - chart_width) / 2
    chart_top = (slide_height - chart_height) / 2 + 0.3  # Slightly below center
    
    if os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            Inches(chart_left),
            Inches(chart_top),
            width=Inches(chart_width)
        )
        print(f"   ✓ Pattern 3: Chart centered on slide")

# ============================================================================
# PATTERN 4: TWO CHARTS SIDE BY SIDE
# ============================================================================
# Use when: Comparing two charts on one slide

def pattern_4_side_by_side(prs, layout_idx, title_text, chart1_path, chart2_path):
    """Two charts side by side"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Add title
    try:
        title = slide.placeholders[0]
        title.text = title_text
        title_bottom = title.top.inches + title.height.inches
        chart_top = title_bottom + 0.3
    except:
        chart_top = 1.5
    
    # Calculate positions for two charts
    margin = 0.5
    gap = 0.5
    chart_width = (10 - 2*margin - gap) / 2  # Two charts with gap
    
    # Chart 1 (left)
    if os.path.exists(chart1_path):
        slide.shapes.add_picture(
            chart1_path,
            Inches(margin),
            Inches(chart_top),
            width=Inches(chart_width)
        )
        print(f"   ✓ Pattern 4: Left chart added")
    
    # Chart 2 (right)
    if os.path.exists(chart2_path):
        slide.shapes.add_picture(
            chart2_path,
            Inches(margin + chart_width + gap),
            Inches(chart_top),
            width=Inches(chart_width)
        )
        print(f"   ✓ Pattern 4: Right chart added")

# ============================================================================
# PATTERN 5: CHART WITH CUSTOM TEXT BOXES
# ============================================================================
# Use when: You need text in specific positions around the chart

def pattern_5_chart_with_annotations(prs, layout_idx, title_text, chart_path):
    """Chart with custom text annotations"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_box.text_frame.text = title_text
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # Chart (left side)
    chart_left = 0.5
    chart_top = 1.2
    chart_width = 5.5
    
    if os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            Inches(chart_left),
            Inches(chart_top),
            width=Inches(chart_width)
        )
    
    # Text box (right side)
    text_box = slide.shapes.add_textbox(
        Inches(6.5),  # Right of chart
        Inches(1.5),
        Inches(3),
        Inches(3.5)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Add bullet points
    p = text_frame.paragraphs[0]
    p.text = "Key Insights:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(12)
    
    insights = [
        "Growth of 15% year-over-year",
        "Mobile viewership increasing",
        "Peak engagement in Q4"
    ]
    
    for insight in insights:
        p = text_frame.add_paragraph()
        p.text = f"• {insight}"
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    print(f"   ✓ Pattern 5: Chart with text annotations")

# ============================================================================
# PATTERN 6: FULL-BLEED CHART (EDGE TO EDGE)
# ============================================================================
# Use when: You want chart to fill entire slide (presentation style)

def pattern_6_full_bleed(prs, layout_idx, chart_path):
    """Chart fills entire slide edge-to-edge"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Fill entire slide
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    if os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            Inches(0),
            Inches(0),
            width=Inches(slide_width),
            height=Inches(slide_height)
        )
        print(f"   ✓ Pattern 6: Full-bleed chart")

# ============================================================================
# PATTERN 7: FIND AND USE TEMPLATE PLACEHOLDERS AUTOMATICALLY
# ============================================================================

def pattern_7_auto_detect(prs, layout_idx, title_text, chart_path):
    """Automatically detect content area and fill it"""
    
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    layout = prs.slide_layouts[layout_idx]
    
    # Add title to first placeholder
    if len(slide.placeholders) > 0:
        slide.placeholders[0].text = title_text
    
    # Find largest placeholder (usually content area)
    largest_ph = None
    largest_area = 0
    
    for ph in layout.placeholders:
        if hasattr(ph, 'width') and hasattr(ph, 'height'):
            area = ph.width * ph.height
            if area > largest_area:
                largest_area = area
                largest_ph = ph
    
    if largest_ph and largest_ph != layout.placeholders[0]:
        # Use the largest placeholder's area
        left = largest_ph.left.inches if largest_ph.left > 0 else 1
        top = largest_ph.top.inches if largest_ph.top > 0 else 1.5
        width = largest_ph.width.inches
        
        # Remove placeholder
        try:
            sp = largest_ph.element
            sp.getparent().remove(sp)
        except:
            pass
        
        # Add chart
        if os.path.exists(chart_path):
            slide.shapes.add_picture(
                chart_path,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
            print(f"   ✓ Pattern 7: Auto-detected content area ({width:.1f}\" wide)")
    else:
        # Fallback to standard position
        pattern_2_below_title(prs, layout_idx, title_text, chart_path)

# ============================================================================
# USAGE EXAMPLES
# ============================================================================

def demonstrate_all_patterns():
    """
    Create examples of all patterns
    """
    
    # Load your template
    TEMPLATE_PATH = "/dbfs/FileStore/templates/your_template.pptx"
    
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template not found: {TEMPLATE_PATH}")
        print("   Using blank presentation for demo")
        prs = Presentation()
    else:
        prs = Presentation(TEMPLATE_PATH)
        print(f"✅ Template loaded")
    
    # Chart paths (from R script)
    chart1 = "/dbfs/FileStore/charts/sports_pie.png"
    chart2 = "/dbfs/FileStore/charts/device_lollipop.png"
    
    print("\nCreating example slides with different patterns...")
    print()
    
    # Pattern 1: Fill placeholder
    if len(prs.slide_layouts) > 1:
        pattern_1_fill_placeholder(prs, 1, "Pattern 1: Fill Placeholder", chart1)
    
    # Pattern 2: Below title (most common)
    if len(prs.slide_layouts) > 1:
        pattern_2_below_title(prs, 1, "Pattern 2: Below Title", chart1)
    
    # Pattern 3: Centered
    if len(prs.slide_layouts) > 5:
        pattern_3_centered(prs, 6, "Pattern 3: Centered", chart1)
    
    # Pattern 4: Side by side
    if len(prs.slide_layouts) > 1:
        pattern_4_side_by_side(prs, 1, "Pattern 4: Side by Side", chart1, chart2)
    
    # Pattern 5: With annotations
    if len(prs.slide_layouts) > 5:
        pattern_5_chart_with_annotations(prs, 6, "Pattern 5: With Annotations", chart1)
    
    # Pattern 6: Full bleed
    if len(prs.slide_layouts) > 5:
        pattern_6_full_bleed(prs, 6, chart1)
    
    # Pattern 7: Auto-detect
    if len(prs.slide_layouts) > 1:
        pattern_7_auto_detect(prs, 1, "Pattern 7: Auto-Detect", chart1)
    
    # Save
    output_path = "/dbfs/FileStore/presentations/positioning_patterns_demo.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    
    print()
    print(f"✅ Demo presentation created: {output_path}")
    print("   Review this file to see all positioning patterns!")

# ============================================================================
# QUICK REFERENCE GUIDE
# ============================================================================

def print_quick_reference():
    """Print a quick reference for common positions"""
    
    print("\n" + "=" * 80)
    print("QUICK REFERENCE: Common Chart Positions (16:9 slide)")
    print("=" * 80)
    print()
    
    print("📍 STANDARD POSITIONS:")
    print("-" * 80)
    print()
    
    print("Full-width chart below title:")
    print("  slide.shapes.add_picture(chart_path,")
    print("                          Inches(1.0),    # left")
    print("                          Inches(1.8),    # top")
    print("                          width=Inches(8.0))")
    print()
    
    print("Centered chart:")
    print("  slide.shapes.add_picture(chart_path,")
    print("                          Inches(1.0),    # left")
    print("                          Inches(2.0),    # top")
    print("                          width=Inches(8.0))")
    print()
    
    print("Wide chart (maximum width with margins):")
    print("  slide.shapes.add_picture(chart_path,")
    print("                          Inches(0.5),    # left")
    print("                          Inches(1.5),    # top")
    print("                          width=Inches(9.0))")
    print()
    
    print("Left half (for side-by-side):")
    print("  slide.shapes.add_picture(chart_path,")
    print("                          Inches(0.5),    # left")
    print("                          Inches(1.5),    # top")
    print("                          width=Inches(4.5))")
    print()
    
    print("Right half (for side-by-side):")
    print("  slide.shapes.add_picture(chart_path,")
    print("                          Inches(5.0),    # left")
    print("                          Inches(1.5),    # top")
    print("                          width=Inches(4.5))")
    print()
    
    print("📍 CALCULATING POSITIONS:")
    print("-" * 80)
    print()
    
    print("To center horizontally:")
    print("  chart_width = 8.0")
    print("  slide_width = 10.0  # Standard 16:9")
    print("  left = (slide_width - chart_width) / 2  # = 1.0")
    print()
    
    print("To position below title:")
    print("  title = slide.placeholders[0]")
    print("  title_bottom = title.top.inches + title.height.inches")
    print("  chart_top = title_bottom + 0.3  # 0.3\" margin")
    print()
    
    print("To use placeholder area:")
    print("  content = slide.placeholders[1]")
    print("  left = content.left.inches")
    print("  top = content.top.inches")
    print("  width = content.width.inches")
    print()
    
    print("=" * 80)

# ============================================================================
# RUN THE SCRIPT
# ============================================================================

if __name__ == "__main__":
    print("=" * 80)
    print("PRACTICAL POSITIONING PATTERNS FOR POWERPOINT")
    print("=" * 80)
    print()
    
    # Show quick reference
    print_quick_reference()
    
    # Create demo presentation
    print("\nCreating demo presentation with all patterns...")
    demonstrate_all_patterns()
    
    print("\n" + "=" * 80)
    print("DONE!")
    print("=" * 80)
    print("\n📚 How to use:")
    print("1. Review the demo presentation")
    print("2. Pick the pattern that matches your needs")
    print("3. Copy that pattern's code into your script")
    print("4. Adjust the coordinates if needed")
    print("=" * 80)
