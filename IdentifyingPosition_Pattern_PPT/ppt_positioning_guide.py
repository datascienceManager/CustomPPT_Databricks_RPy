"""
============================================================================
POSITIONING GUIDE: Charts & Text in Custom PowerPoint Templates
============================================================================
This script helps you:
1. Find exact placeholder positions in your template
2. Position ggplot charts precisely
3. Add text in specific locations
4. Avoid overlapping elements
============================================================================
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

print("=" * 80)
print("POWERPOINT POSITIONING DIAGNOSTIC TOOL")
print("=" * 80)
print()

# ============================================================================
# STEP 1: ANALYZE YOUR TEMPLATE
# ============================================================================

def analyze_template(template_path):
    """
    Comprehensive template analysis showing:
    - All layouts and their names
    - All placeholders and their positions
    - Exact coordinates for each element
    """
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    prs = Presentation(template_path)
    print(f"✅ Template loaded: {template_path}")
    print(f"   Slide size: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
    print()
    
    # Analyze each layout
    for layout_idx, layout in enumerate(prs.slide_layouts):
        print("=" * 80)
        print(f"LAYOUT {layout_idx}: {layout.name}")
        print("=" * 80)
        
        if len(layout.placeholders) == 0:
            print("   ⚠ No placeholders (blank layout)")
            print()
            continue
        
        # Show each placeholder
        for ph_idx, placeholder in enumerate(layout.placeholders):
            print(f"\n📍 Placeholder [{ph_idx}]: {placeholder.name}")
            print(f"   Type: {placeholder.placeholder_format.type}")
            
            # Position information
            left_inches = placeholder.left.inches if placeholder.left > 0 else 0
            top_inches = placeholder.top.inches if placeholder.top > 0 else 0
            width_inches = placeholder.width.inches
            height_inches = placeholder.height.inches
            
            print(f"   Position (inches):")
            print(f"      Left:   {left_inches:.2f}\"")
            print(f"      Top:    {top_inches:.2f}\"")
            print(f"      Width:  {width_inches:.2f}\"")
            print(f"      Height: {height_inches:.2f}\"")
            
            # Calculate right and bottom edges
            right = left_inches + width_inches
            bottom = top_inches + height_inches
            print(f"   Boundaries:")
            print(f"      Right edge:  {right:.2f}\"")
            print(f"      Bottom edge: {bottom:.2f}\"")
            
            # Show safe area for content
            print(f"   💡 Safe area for images (inside this placeholder):")
            print(f"      add_image_to_slide(slide, chart_path,")
            print(f"                        left={left_inches:.2f}, top={top_inches:.2f},")
            print(f"                        width={width_inches:.2f}, height={height_inches:.2f})")
        
        print()
    
    # Analyze actual slide if template has examples
    if len(prs.slides) > 0:
        print("=" * 80)
        print("EXAMPLE SLIDE ANALYSIS (First slide in template)")
        print("=" * 80)
        slide = prs.slides[0]
        
        for idx, shape in enumerate(slide.shapes):
            print(f"\nShape {idx}: {shape.name}")
            print(f"   Type: {shape.shape_type}")
            
            if hasattr(shape, 'left'):
                print(f"   Position:")
                print(f"      Left:   {shape.left.inches:.2f}\"")
                print(f"      Top:    {shape.top.inches:.2f}\"")
                print(f"      Width:  {shape.width.inches:.2f}\"")
                print(f"      Height: {shape.height.inches:.2f}\"")

# ============================================================================
# STEP 2: HELPER FUNCTIONS FOR PRECISE POSITIONING
# ============================================================================

def calculate_centered_position(slide_width, slide_height, image_width, image_height):
    """Calculate position to center an image on slide"""
    left = (slide_width - image_width) / 2
    top = (slide_height - image_height) / 2
    return left, top

def calculate_chart_position_in_placeholder(placeholder, chart_aspect_ratio=1.67):
    """
    Calculate best position to fit chart inside a placeholder
    chart_aspect_ratio: width/height ratio of your chart (e.g., 10/6 = 1.67)
    """
    ph_left = placeholder.left.inches if placeholder.left > 0 else 0
    ph_top = placeholder.top.inches if placeholder.top > 0 else 0
    ph_width = placeholder.width.inches
    ph_height = placeholder.height.inches
    
    # Calculate chart size to fit in placeholder
    # Try fitting by width first
    chart_width = ph_width
    chart_height = chart_width / chart_aspect_ratio
    
    # If height is too tall, fit by height instead
    if chart_height > ph_height:
        chart_height = ph_height
        chart_width = chart_height * chart_aspect_ratio
    
    # Center the chart in placeholder
    left = ph_left + (ph_width - chart_width) / 2
    top = ph_top + (ph_height - chart_height) / 2
    
    return {
        'left': left,
        'top': top,
        'width': chart_width,
        'height': chart_height
    }

def add_image_to_placeholder_area(slide, placeholder_idx, image_path, margin=0.2):
    """
    Add image inside a placeholder area with margins
    margin: inches to leave as padding around image
    """
    try:
        placeholder = slide.placeholders[placeholder_idx]
        
        # Get placeholder dimensions
        ph_left = placeholder.left.inches if placeholder.left > 0 else 0
        ph_top = placeholder.top.inches if placeholder.top > 0 else 0
        ph_width = placeholder.width.inches
        ph_height = placeholder.height.inches
        
        # Calculate position with margins
        left = ph_left + margin
        top = ph_top + margin
        width = ph_width - (2 * margin)
        height = ph_height - (2 * margin)
        
        # Add image
        if os.path.exists(image_path):
            # Remove placeholder first (optional)
            sp = placeholder.element
            sp.getparent().remove(sp)
            
            # Add image
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
            print(f"   ✓ Image added to placeholder {placeholder_idx}")
            return True
        else:
            print(f"   ✗ Image not found: {image_path}")
            return False
            
    except (KeyError, IndexError) as e:
        print(f"   ✗ Placeholder {placeholder_idx} not found: {e}")
        return False

def add_chart_below_title(slide, image_path, title_placeholder_idx=0, 
                         margin_top=0.3, chart_width=8, chart_height=4):
    """
    Add chart below the title placeholder
    """
    try:
        title_ph = slide.placeholders[title_placeholder_idx]
        
        # Calculate position below title
        title_bottom = title_ph.top.inches + title_ph.height.inches
        chart_left = (10 - chart_width) / 2  # Center horizontally (assuming 16:9 slide)
        chart_top = title_bottom + margin_top
        
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(chart_left),
                Inches(chart_top),
                width=Inches(chart_width)
            )
            print(f"   ✓ Chart added below title")
            return True
        else:
            print(f"   ✗ Image not found: {image_path}")
            return False
            
    except Exception as e:
        print(f"   ✗ Error adding chart: {e}")
        return False

def create_custom_text_box(slide, text, left, top, width, height, 
                           font_size=16, bold=False, color=None, 
                           align=PP_ALIGN.LEFT):
    """
    Create a text box at exact position
    """
    text_box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.alignment = align
        if color:
            paragraph.font.color.rgb = color
    
    return text_box

# ============================================================================
# STEP 3: PRACTICAL EXAMPLES FOR YOUR TEMPLATE
# ============================================================================

def example_positioning_strategies(template_path):
    """
    Show different strategies for positioning with your template
    """
    
    prs = Presentation(template_path)
    
    print("\n" + "=" * 80)
    print("POSITIONING STRATEGIES FOR YOUR TEMPLATE")
    print("=" * 80)
    
    # Strategy 1: Use placeholder dimensions
    print("\n📋 STRATEGY 1: Fill placeholder area")
    print("-" * 80)
    print("Use when: You want chart to fit exactly in a placeholder")
    print()
    print("Code:")
    print("""
# Get placeholder dimensions
placeholder = slide.placeholders[1]  # Content placeholder
ph_left = placeholder.left.inches
ph_top = placeholder.top.inches
ph_width = placeholder.width.inches
ph_height = placeholder.height.inches

# Remove placeholder
sp = placeholder.element
sp.getparent().remove(sp)

# Add chart in exact same space
slide.shapes.add_picture(
    "/dbfs/FileStore/charts/your_chart.png",
    Inches(ph_left),
    Inches(ph_top),
    width=Inches(ph_width)  # Height auto-scales
)
    """)
    
    # Strategy 2: Below title
    print("\n📋 STRATEGY 2: Position below title")
    print("-" * 80)
    print("Use when: Title at top, chart below")
    print()
    print("Code:")
    print("""
# Find where title ends
title = slide.placeholders[0]
title_bottom = title.top.inches + title.height.inches

# Add chart below with margin
chart_top = title_bottom + 0.3  # 0.3 inch margin
chart_left = 1.0  # 1 inch from left
chart_width = 8.0

slide.shapes.add_picture(
    chart_path,
    Inches(chart_left),
    Inches(chart_top),
    width=Inches(chart_width)
)
    """)
    
    # Strategy 3: Manual coordinates
    print("\n📋 STRATEGY 3: Manual coordinates")
    print("-" * 80)
    print("Use when: You know exactly where you want it")
    print()
    print("Slide dimensions:")
    print(f"   Width:  {prs.slide_width.inches:.2f}\"")
    print(f"   Height: {prs.slide_height.inches:.2f}\"")
    print()
    print("Code:")
    print("""
# Manual positioning (16:9 slide)
slide.shapes.add_picture(
    chart_path,
    Inches(1.0),    # 1 inch from left
    Inches(2.0),    # 2 inches from top
    width=Inches(8.0)  # 8 inches wide
)

# For centering:
chart_width = 8.0
slide_width = 10.0
left = (slide_width - chart_width) / 2  # = 1.0

slide.shapes.add_picture(
    chart_path,
    Inches(left),
    Inches(2.0),
    width=Inches(chart_width)
)
    """)
    
    # Strategy 4: Grid-based
    print("\n📋 STRATEGY 4: Grid-based layout")
    print("-" * 80)
    print("Use when: Multiple charts on one slide")
    print()
    print("Code:")
    print("""
# 2x2 grid of charts
margin = 0.5
chart_width = (10 - 3*margin) / 2  # Two charts with margins
chart_height = 2.5

positions = [
    (margin, 1.0),                    # Top-left
    (margin + chart_width + margin, 1.0),  # Top-right
    (margin, 1.0 + chart_height + margin), # Bottom-left
    (margin + chart_width + margin, 1.0 + chart_height + margin)  # Bottom-right
]

for idx, (left, top) in enumerate(positions):
    slide.shapes.add_picture(
        f"/dbfs/FileStore/charts/chart_{idx}.png",
        Inches(left),
        Inches(top),
        width=Inches(chart_width)
    )
    """)

# ============================================================================
# STEP 4: VISUAL COORDINATE SYSTEM
# ============================================================================

def show_coordinate_system():
    """Show PowerPoint coordinate system"""
    
    print("\n" + "=" * 80)
    print("POWERPOINT COORDINATE SYSTEM")
    print("=" * 80)
    print()
    print("Standard 16:9 slide (10\" × 5.625\"):")
    print()
    print("  (0,0) ──────────────────────────── (10,0)")
    print("    │                                    │")
    print("    │         YOUR CONTENT               │")
    print("    │                                    │")
    print("    │                                    │")
    print("  (0,5.625) ────────────────────── (10,5.625)")
    print()
    print("Common positions:")
    print("  • Top-left corner:     (0, 0)")
    print("  • Top-center:          (5, 0)")
    print("  • Center:              (5, 2.8)")
    print("  • Bottom-left:         (0, 5.625)")
    print()
    print("Safe content area (with 0.5\" margins):")
    print("  • Left:   0.5\"")
    print("  • Right:  9.5\"")
    print("  • Top:    0.5\"")
    print("  • Bottom: 5.125\"")
    print()
    print("Typical chart positions:")
    print("  • Full-width chart:")
    print("    left=1.0, top=1.5, width=8.0")
    print()
    print("  • Chart below title:")
    print("    left=1.0, top=1.8, width=8.0")
    print()
    print("  • Side-by-side charts:")
    print("    Chart 1: left=0.5, top=1.5, width=4.5")
    print("    Chart 2: left=5.0, top=1.5, width=4.5")

# ============================================================================
# STEP 5: DEBUGGING TOOL
# ============================================================================

def create_positioning_test_slide(template_path, output_path):
    """
    Create a test slide showing grid overlay to help with positioning
    """
    
    prs = Presentation(template_path)
    
    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Usually blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Draw grid lines
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # Vertical lines every inch
    for i in range(int(slide_width) + 1):
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(i), Inches(0),
            Inches(0.01), Inches(slide_height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        # Add label
        label = slide.shapes.add_textbox(
            Inches(i + 0.1), Inches(0.1),
            Inches(0.5), Inches(0.3)
        )
        label.text_frame.text = f"{i}\""
        label.text_frame.paragraphs[0].font.size = Pt(10)
    
    # Horizontal lines every inch
    for i in range(int(slide_height) + 1):
        line = slide.shapes.add_shape(
            1,  # Line
            Inches(0), Inches(i),
            Inches(slide_width), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        # Add label
        label = slide.shapes.add_textbox(
            Inches(0.1), Inches(i + 0.1),
            Inches(0.5), Inches(0.3)
        )
        label.text_frame.text = f"{i}\""
        label.text_frame.paragraphs[0].font.size = Pt(10)
    
    # Add title
    title = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(0.5)
    )
    title.text_frame.text = "Positioning Grid (1 inch squares)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    prs.save(output_path)
    print(f"\n✅ Test slide created: {output_path}")
    print("   Use this to visually measure positions in your template")

# ============================================================================
# MAIN EXECUTION
# ============================================================================

if __name__ == "__main__":
    
    # Set your template path
    TEMPLATE_PATH = "/dbfs/FileStore/templates/your_template.pptx"
    
    print("Choose an option:")
    print("1. Analyze template (find placeholder positions)")
    print("2. Show positioning strategies")
    print("3. Show coordinate system")
    print("4. Create grid test slide")
    print()
    
    # For Databricks, run all options
    print("Running all diagnostic tools...\n")
    
    # Option 1: Analyze template
    if os.path.exists(TEMPLATE_PATH):
        analyze_template(TEMPLATE_PATH)
    else:
        print(f"❌ Template not found: {TEMPLATE_PATH}")
        print("   Please upload your template and update TEMPLATE_PATH")
    
    # Option 2: Show strategies
    if os.path.exists(TEMPLATE_PATH):
        example_positioning_strategies(TEMPLATE_PATH)
    
    # Option 3: Show coordinates
    show_coordinate_system()
    
    # Option 4: Create test slide
    if os.path.exists(TEMPLATE_PATH):
        test_output = "/dbfs/FileStore/presentations/positioning_test.pptx"
        create_positioning_test_slide(TEMPLATE_PATH, test_output)
    
    print("\n" + "=" * 80)
    print("DIAGNOSTIC COMPLETE")
    print("=" * 80)
    print("\nNext steps:")
    print("1. Review the placeholder positions above")
    print("2. Update your script with correct coordinates")
    print("3. Use the test slide to verify positions visually")
    print("=" * 80)
